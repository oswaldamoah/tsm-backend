"""
Turn a presentation spec (from the agent) into a real .pptx file.

Charts are added as native PowerPoint charts rather than images, so the numbers
stay live and editable when someone opens the deck.
"""

from io import BytesIO
from typing import Optional

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.util import Inches, Pt


SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

ACCENT = RGBColor(0x1F, 0x6F, 0xEB)
INK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x66, 0x66, 0x70)

_CHART_TYPES = {
    "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "stacked_bar": XL_CHART_TYPE.COLUMN_STACKED,
    "line": XL_CHART_TYPE.LINE_MARKERS,
    "area": XL_CHART_TYPE.AREA,
    "pie": XL_CHART_TYPE.PIE,
    "doughnut": XL_CHART_TYPE.DOUGHNUT,
}


def _add_chart(slide, chart_spec: dict, left, top, width, height) -> None:
    chart_data = CategoryChartData()
    chart_data.categories = chart_spec.get("categories") or []
    for series in chart_spec.get("series") or []:
        chart_data.add_series(series.get("name") or "Value", series.get("values") or [])

    chart_type = _CHART_TYPES.get(chart_spec.get("type", "bar"), XL_CHART_TYPE.COLUMN_CLUSTERED)
    frame = slide.shapes.add_chart(chart_type, left, top, width, height, chart_data)
    chart = frame.chart

    multi_series = len(chart_spec.get("series") or []) > 1
    is_circular = chart_spec.get("type") in ("pie", "doughnut")
    chart.has_legend = multi_series or is_circular
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False

    if chart.has_title:
        chart.has_title = False


def _add_bullets(slide, bullets: list[str], left, top, width, height, size: int = 16) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True

    for index, bullet in enumerate(bullets):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = f"•  {bullet}"
        paragraph.space_after = Pt(10)
        for run in paragraph.runs:
            run.font.size = Pt(size)
            run.font.color.rgb = INK


def _title_only_slide(prs: Presentation, title: str):
    """Blank layout with our own title textbox - avoids template placeholder surprises."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), SLIDE_WIDTH - Inches(1.4), Inches(0.9))
    paragraph = box.text_frame.paragraphs[0]
    paragraph.text = title
    run = paragraph.runs[0]
    run.font.size = Pt(30)
    run.font.bold = True
    run.font.color.rgb = INK
    return slide


def build_pptx(presentation: dict, charts: Optional[list[dict]] = None) -> BytesIO:
    """Render the deck and return it as an in-memory .pptx stream."""
    charts = charts or []

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # --- Title slide ---
    cover = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = cover.shapes.add_textbox(Inches(0.9), Inches(2.4), SLIDE_WIDTH - Inches(1.8), Inches(1.4))
    title_paragraph = title_box.text_frame.paragraphs[0]
    title_paragraph.text = presentation.get("title") or "Presentation"
    title_run = title_paragraph.runs[0]
    title_run.font.size = Pt(44)
    title_run.font.bold = True
    title_run.font.color.rgb = ACCENT

    subtitle = presentation.get("subtitle")
    if subtitle:
        subtitle_box = cover.shapes.add_textbox(Inches(0.9), Inches(3.8), SLIDE_WIDTH - Inches(1.8), Inches(0.8))
        subtitle_paragraph = subtitle_box.text_frame.paragraphs[0]
        subtitle_paragraph.text = subtitle
        subtitle_run = subtitle_paragraph.runs[0]
        subtitle_run.font.size = Pt(20)
        subtitle_run.font.color.rgb = MUTED

    # --- Content slides ---
    for slide_spec in presentation.get("slides") or []:
        slide = _title_only_slide(prs, slide_spec.get("title") or "Slide")
        bullets = slide_spec.get("bullets") or []
        chart_index = slide_spec.get("chartIndex")
        chart_spec = charts[chart_index] if isinstance(chart_index, int) and 0 <= chart_index < len(charts) else None

        if chart_spec and bullets:
            _add_bullets(slide, bullets, Inches(0.7), Inches(1.6), Inches(5.2), Inches(5.0))
            _add_chart(slide, chart_spec, Inches(6.3), Inches(1.6), Inches(6.3), Inches(5.0))
        elif chart_spec:
            _add_chart(slide, chart_spec, Inches(1.2), Inches(1.6), SLIDE_WIDTH - Inches(2.4), Inches(5.0))
        else:
            _add_bullets(slide, bullets, Inches(0.7), Inches(1.7), SLIDE_WIDTH - Inches(1.4), Inches(4.8), size=18)

        notes = slide_spec.get("notes")
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    stream = BytesIO()
    prs.save(stream)
    stream.seek(0)
    return stream


def safe_filename(title: str) -> str:
    """A title turned into something safe for a Content-Disposition filename."""
    cleaned = "".join(ch if (ch.isalnum() or ch in " -_") else "" for ch in (title or "presentation"))
    cleaned = "-".join(cleaned.split()) or "presentation"
    return f"{cleaned[:60]}.pptx"
